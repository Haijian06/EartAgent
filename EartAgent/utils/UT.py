# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
from EartAgent.Agent.text_agents import *
import logging
from datetime import datetime
import serpapi
import docx
from docx import Document
import PyPDF2
from pptx import Presentation
from langchain_community.document_loaders import AsyncHtmlLoader
from langchain_community.document_transformers import Html2TextTransformer
import re
import os
import json
import mammoth
import textract


class UtilityTools:
    def __init__(self):
        logging.basicConfig(level=logging.INFO)
        self.init_log_dir()
        self.log_file_path = self.init_log_file()
        self.initialized = False  # Initialization flag to prevent duplicate logging

    def init_log_dir(self):
        # Create a log directory in the 'log_time' folder within the current working directory
        self.log_dir = os.path.join(os.getcwd(), 'log_time')
        if not os.path.exists(self.log_dir):
            os.makedirs(self.log_dir)

    def init_log_file(self):
        # Create a base log file for appending content later
        timestamp = datetime.now().strftime("%Y-%m-%d %H-%M-%S")
        log_file_path = os.path.join(self.log_dir, f"log_{timestamp}.json")
        return log_file_path

    def log_agent(self, agent_instance):
        # Log only during the first instantiation
        if not self.initialized:
            timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            log_data = {
                "timestamp": timestamp,
                "agent_config": {
                    "name": agent_instance.config.name,
                    "model_name": agent_instance.config.model_name,
                    "temperature": agent_instance.config.temperature,
                    "max_tokens": agent_instance.config.max_tokens,
                    "remember": agent_instance.config.remember
                }
            }
            with open(self.log_file_path, 'a') as file:
                if os.stat(self.log_file_path).st_size == 0:
                    file.write("[\n")
                else:
                    file.seek(file.tell() - 2, os.SEEK_SET)
                    file.write(",\n")
                json.dump(log_data, file, indent=4)
                file.write("\n]")
            # Set log format and level
            logging.basicConfig(format='%(message)s', level=logging.INFO)

            # ANSI escape code - Blue
            BLUE = '\033[94m'
            ENDC = '\033[0m'
            # Set httpx log level to WARNING or higher to avoid logging every request and response
            logging.getLogger("httpx").setLevel(logging.WARNING)

            logging.info(f"{BLUE}Logged agent initialization data at {timestamp}{ENDC}")
            self.initialized = True  # Ensure logging happens only once

    def get_date(self):
        """
        Get the current date and time
        """
        return datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    def serpapi_search(self, query, api_Key):
        """
        Perform a search using SerpAPI (default is Google)
        """
        print(">>>>>>>>>>>performing serpapi_search")
        params = {
            "engine": "google",
            "q": query,
            "api_key": api_Key
        }

        search = serpapi.search(params)
        dict_search = dict(search)
        text_collection_ready_to_LLM = ''
        dict_collection = {}

        try:
            answer_box = dict_search['answer_box']
            answer_text = answer_box['title']
            answer_link = answer_box['link']
            text_collection_ready_to_LLM += answer_text
            dict_collection[answer_text] = answer_link
        except KeyError:
            print("No answer box in the response, returning organic_results only")

        organic_results = dict_search.get('organic_results', [])
        for result in organic_results:
            result_title = result['title']
            result_link = result['link']
            text_collection_ready_to_LLM += " " + result_title
            dict_collection[result_title] = result_link

        print(">>>>>>>>>>>serpapi_search completed")
        return text_collection_ready_to_LLM, dict_collection

    def read_docx(self, file_path):
        """
        Read content from a docx file
        :param file_path: Path to the docx file
        :return: List of file content, each element corresponds to a paragraph
        """
        doc = docx.Document(file_path)
        text = [para.text for para in doc.paragraphs]
        return text

    def read_doc(self, file_path):
        """
        Read content from a doc file
        :param file_path: Path to the doc file
        :return: List of file content, each element corresponds to a line
        """
        # Read document content
        text = textract.process(file_path)
        return text.decode('utf-8')

    def read_txt(self, file_path):
        """
        Read content from a txt file
        :param file_path: Path to the txt file
        :return: List of file content, each element corresponds to a line
        """
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read().splitlines()
        return text

    def write_docx(self, file_path, content):
        """
        Write content to a Word document
        """
        doc = Document()
        for line in content:
            doc.add_paragraph(line)
        doc.save(file_path)

    def read_pdf(self, file_path):
        """
        Read content from a PDF file
        """
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = [page.extract_text() for page in pdf_reader.pages]
        return text

    def read_ppt(self, file_path):
        """
        Read content from a PowerPoint file
        """
        ppt = Presentation(file_path)
        text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.append(shape.text)
        return text

    def web_crawler_all(self, url):
        """
        Crawl all information from a webpage
        """
        headers = {}
        loader = AsyncHtmlLoader([url])
        docs = loader.load()
        html2text_transformer = Html2TextTransformer()
        docs_transformed = html2text_transformer.transform_documents(docs)
        print(docs_transformed)
        text = str(docs_transformed)
        clean_text = re.sub(r'[^\w\s]', '', text)
        clean_text = re.sub(r'[n]', '', clean_text)
        return clean_text

    def web_crawler_by_LLM(self, url, prompt, agent):
        """
        Adaptive crawler that leverages LLM for enhanced performance
        """
        loader = AsyncHtmlLoader([url])
        docs = loader.load()
        html2text_transformer = Html2TextTransformer()
        docs_transformed = html2text_transformer.transform_documents(docs)
        text = str(docs_transformed)
        clean_text = re.sub(r'[^\w\s]', '', text)
        clean_text = re.sub(r'[n]', '', clean_text)
        transformed_text = agent(
            f'Process this text and extract relevant content based on the prompt "{prompt}" while summarizing it succinctly within 500 words: {clean_text}')
        return transformed_text

    def search_crawler(self, api_key, q, max_crawler_count, agent):
        """
        Integration of search and crawling
        """
        if isinstance(q, str):
            text_collection_ready_to_LLM, dict_collection = self.serpapi_search(q, api_key)
        else:
            raise ValueError(
                "Don't guess what this sentence means, just enter the question as a string and there is no problem")

        crawler_res = ''
        count = 0
        for key in dict_collection:
            if count < max_crawler_count:
                url = dict_collection[key]
                if url:
                    print('url:', url)
                    crawler_res += self.web_crawler_by_LLM(url, q, agent)
                    count += 1
                else:
                    raise ValueError('Cannot find the URL from dict_collection function')
            else:
                print("The maximum number of iterations for crawling has been reached!")
                break
        return crawler_res

    def reflect_on_response(self, scenario: str, agent_response: str, agent, iterations: int = 1) -> str:
        """
        Let the Agent reflect on its response in a specific scenario, can iterate multiple times for deeper reflection.
        Args:
            scenario (str): Describes the scenario the agent needs to reflect on, e.g., "code explanation" or "answering questions".
            agent_response (str): The agent's original response to the scenario.
            agent (Agent): The agent instance providing reflection capability.
            iterations (int): Number of reflection iterations.

        Returns:
            str: Reflection report, describing the strengths and weaknesses of the original response and suggestions for improvement.
        """
        current_response = agent_response
        reflection_report = ""

        for i in range(iterations):
            # Define a prompt for the scenario, including the current response
            prompt = f"反思以下对情景的回应'{scenario}': '{current_response}'. 考虑回应的清晰度、相关性和准确性。提供改进建议优化回答。"

            # Using Agents to Generate Reflection Reports
            reflection = agent.chat(prompt)
            agent.speak(reflection)
            reflection_report += f"Iteration {i + 1}: {reflection}\n\n"

            # Update the current answer as a reflective suggestion so that the next iteration can be improved further
            current_response = reflection

        return reflection_report

    def some_other_utility_function(self):
        """
        Add other useful utility functions here

        """
        pass



